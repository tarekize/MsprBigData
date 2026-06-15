"""
Génère le support de soutenance PowerPoint — Electio-Analytics POC
Sauvegarde dans outputs/soutenance_electio_analytics.pptx

Dépendance : pip install python-pptx
"""

import os
import sys

try:
    from pptx import Presentation
    from pptx.util import Inches, Pt, Emu
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN
    from pptx.util import Cm
except ImportError:
    print("❌ python-pptx non installé. Lancez : pip install python-pptx")
    sys.exit(1)

_SCRIPT_DIR = os.path.dirname(os.path.abspath(__file__))
_MSPR_FINAL = os.path.abspath(os.path.join(_SCRIPT_DIR, '..', '..'))
OUTPUT_DIR = os.path.join(_MSPR_FINAL, 'outputs')
os.makedirs(OUTPUT_DIR, exist_ok=True)

IMAGES_DIR = OUTPUT_DIR

# Couleurs du thème
BLUE_DARK = RGBColor(0x1A, 0x52, 0x76)    # Bleu marine
BLUE_MID = RGBColor(0x29, 0x80, 0xB9)     # Bleu moyen
BLUE_LIGHT = RGBColor(0xD6, 0xEA, 0xF8)   # Bleu clair
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
ORANGE = RGBColor(0xF3, 0x9C, 0x12)
GREEN = RGBColor(0x27, 0xAE, 0x60)
RED = RGBColor(0xE7, 0x4C, 0x3C)
GREY = RGBColor(0x7F, 0x8C, 0x8D)

prs = Presentation()
prs.slide_width = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK_LAYOUT = prs.slide_layouts[6]  # Entièrement vide


def add_bg_rect(slide, color, left=0, top=0, width=None, height=None):
    """Ajoute un rectangle de fond."""
    w = width or prs.slide_width
    h = height or prs.slide_height
    shape = slide.shapes.add_shape(1, left, top, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_text_box(slide, text, left, top, width, height,
                 font_size=18, bold=False, color=WHITE,
                 align=PP_ALIGN.LEFT, italic=False):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox


def add_slide_header(slide, title, subtitle=None, bg_color=BLUE_DARK):
    """Bandeau de titre en haut de chaque slide."""
    add_bg_rect(slide, bg_color, 0, 0, prs.slide_width, Inches(1.4))
    add_text_box(slide, title,
                 Inches(0.4), Inches(0.1), Inches(12), Inches(0.75),
                 font_size=24, bold=True, color=WHITE)
    if subtitle:
        add_text_box(slide, subtitle,
                     Inches(0.4), Inches(0.85), Inches(12), Inches(0.45),
                     font_size=14, color=BLUE_LIGHT)


def add_bullet_box(slide, items, left, top, width, height,
                   title=None, title_color=BLUE_DARK, font_size=14):
    """Boîte avec titre optionnel et liste de bullets."""
    if title:
        add_text_box(slide, title, left, top, width, Inches(0.4),
                     font_size=14, bold=True, color=title_color)
        top += Inches(0.4)
        height -= Inches(0.4)

    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()
        p.space_before = Pt(4)
        run = p.add_run()
        run.text = f"▸  {item}"
        run.font.size = Pt(font_size)
        run.font.color.rgb = RGBColor(0x2C, 0x3E, 0x50)


def add_metric_card(slide, value, label, left, top, width=Inches(2.2), height=Inches(1.3),
                    bg=BLUE_MID, val_color=WHITE, label_color=BLUE_LIGHT):
    card = slide.shapes.add_shape(1, left, top, width, height)
    card.fill.solid()
    card.fill.fore_color.rgb = bg
    card.line.color.rgb = WHITE

    add_text_box(slide, value, left, top + Inches(0.1), width, Inches(0.7),
                 font_size=28, bold=True, color=val_color, align=PP_ALIGN.CENTER)
    add_text_box(slide, label, left, top + Inches(0.75), width, Inches(0.5),
                 font_size=11, color=label_color, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — PAGE DE TITRE
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK_LAYOUT)
add_bg_rect(slide, BLUE_DARK)

# Bandeau orange en bas
add_bg_rect(slide, ORANGE, 0, Inches(6.8), prs.slide_width, Inches(0.7))

add_text_box(slide, "ELECTIO-ANALYTICS",
             Inches(1), Inches(1.0), Inches(11), Inches(0.6),
             font_size=16, bold=False, color=ORANGE, align=PP_ALIGN.CENTER)

add_text_box(slide, "Prédiction des Résultats Électoraux\nen Nouvelle-Aquitaine",
             Inches(1), Inches(1.6), Inches(11), Inches(1.8),
             font_size=36, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

add_text_box(slide, "Preuve de Concept (POC) — MSPR TPRE813",
             Inches(1), Inches(3.5), Inches(11), Inches(0.5),
             font_size=18, italic=True, color=BLUE_LIGHT, align=PP_ALIGN.CENTER)

add_text_box(slide, "Modèle XGBoost · 83% Accuracy · 127 260 enregistrements · 12 départements",
             Inches(1), Inches(4.2), Inches(11), Inches(0.5),
             font_size=13, color=BLUE_LIGHT, align=PP_ALIGN.CENTER)

add_text_box(slide, "Équipe Data Science | Juin 2026",
             Inches(1), Inches(6.85), Inches(11), Inches(0.5),
             font_size=12, color=WHITE, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — CONTEXTE ET OBJECTIFS
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK_LAYOUT)
add_bg_rect(slide, RGBColor(0xF8, 0xF9, 0xFA))
add_slide_header(slide, "Contexte & Objectifs",
                 "Electio-Analytics — Validation de l'approche prédictive")

add_bullet_box(slide,
    ["La start-up Electio-Analytics souhaite prédire les résultats électoraux "
     "à partir d'indicateurs socio-économiques publics",
     "Périmètre : Nouvelle-Aquitaine — 12 départements, 127 260 enregistrements communaux",
     "Données sources : INSEE 2016/2020, résultats élections 2012 & 2017, indicateurs sécurité",
     "Livrable : pipeline ETL + modèle ML + dashboard de visualisation",
     "Objectif d'accuracy : ≥ 80% — Résultat obtenu : 83,07%"],
    Inches(0.5), Inches(1.6), Inches(7.5), Inches(4.5),
    title="Contexte du projet", font_size=13)

# Boîte résumé résultat
add_bg_rect(slide, BLUE_DARK, Inches(8.3), Inches(1.6), Inches(4.5), Inches(4.5))
add_text_box(slide, "Résultat clé", Inches(8.5), Inches(1.7), Inches(4.1), Inches(0.4),
             font_size=13, bold=True, color=ORANGE)
add_text_box(slide,
    "Prédiction 2022 :\nNOUVELLE-AQUITAINE → MACRON\n\nRésultat réel :\nMACRON ✅\n\nConfiance : 88,58%",
    Inches(8.5), Inches(2.2), Inches(4.1), Inches(3.5),
    font_size=16, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — CHOIX GÉOGRAPHIQUE
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK_LAYOUT)
add_bg_rect(slide, RGBColor(0xF8, 0xF9, 0xFA))
add_slide_header(slide, "Justification du Choix Géographique",
                 "Nouvelle-Aquitaine — Critères de sélection")

criteres = [
    ("Disponibilité des données", "INSEE fournit des données exhaustives à l'échelle communale pour cette région"),
    ("Représentativité politique", "11 dép. Macron + 1 dép. Le Pen → diversité pour entraîner le modèle"),
    ("Volumétrie exploitable", "127 260 enregistrements — suffisant sans être excessif pour un POC"),
    ("Historique disponible", "Résultats électoraux 2012 et 2017 accessibles en open data"),
    ("Diversité socio-éco", "Métropole (Bordeaux), zones rurales (Creuse), côtes, agriculture"),
]

for i, (titre, desc) in enumerate(criteres):
    top = Inches(1.6 + i * 1.0)
    add_bg_rect(slide, BLUE_LIGHT if i % 2 == 0 else WHITE,
                Inches(0.3), top, Inches(12.5), Inches(0.85))
    add_text_box(slide, f"✓  {titre} :", Inches(0.5), top + Inches(0.05),
                 Inches(3.5), Inches(0.75), font_size=12, bold=True, color=BLUE_DARK)
    add_text_box(slide, desc, Inches(4.0), top + Inches(0.05),
                 Inches(8.5), Inches(0.75), font_size=12, color=RGBColor(0x2C, 0x3E, 0x50))

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — ARCHITECTURE DU PIPELINE ETL
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK_LAYOUT)
add_bg_rect(slide, RGBColor(0xF8, 0xF9, 0xFA))
add_slide_header(slide, "Pipeline ETL — Architecture des Données",
                 "Flux de données de la collecte à la prédiction")

etapes = [
    ("1. Sources", "Élections 2012/2017\nINSEE 2016/2020\nSécurité RALFSS", BLUE_DARK),
    ("2. Extraction", "run_prep.py\nChargement CSV/XLS/XLSX\nEncodage UTF-8/Latin1", BLUE_MID),
    ("3. Transformation", "Standardisation CODGEO\nCalcul deltas 2016→2020\nNettoyage & normalisation", RGBColor(0x1A, 0x8C, 0x6A)),
    ("4. Chargement", "NoSQL : MongoDB\n+ TinyDB (fallback)\nExport CSV final", RGBColor(0x8E, 0x44, 0xAD)),
    ("5. ML & Prédiction", "XGBoost Classifier\n83,07% accuracy\nExport JSON", ORANGE),
]

for i, (titre, desc, color) in enumerate(etapes):
    left = Inches(0.3 + i * 2.55)
    add_bg_rect(slide, color, left, Inches(1.6), Inches(2.3), Inches(4.5))
    add_text_box(slide, titre, left + Inches(0.1), Inches(1.7), Inches(2.1), Inches(0.5),
                 font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_text_box(slide, desc, left + Inches(0.1), Inches(2.3), Inches(2.1), Inches(3.5),
                 font_size=11, color=WHITE, align=PP_ALIGN.CENTER)
    if i < 4:
        add_text_box(slide, "→", left + Inches(2.3), Inches(3.4), Inches(0.25), Inches(0.5),
                     font_size=22, bold=True, color=BLUE_MID, align=PP_ALIGN.CENTER)

add_text_box(slide,
    "Tables NoSQL : raw_data (données brutes) · final_data (données transformées) · MongoDB / TinyDB fallback",
    Inches(0.3), Inches(6.3), Inches(12.5), Inches(0.5),
    font_size=11, italic=True, color=GREY, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — MODÈLE CONCEPTUEL DE DONNÉES
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK_LAYOUT)
add_bg_rect(slide, RGBColor(0xF8, 0xF9, 0xFA))
add_slide_header(slide, "Modèle Conceptuel de Données (MCD)",
                 "Structure des entités et relations — Electio-Analytics")

entites = [
    ("COMMUNE", ["codgeo (PK)", "nom_commune", "code_dept", "region"], Inches(0.3), Inches(1.7)),
    ("ELECTION", ["id_election (PK)", "codgeo (FK)", "annee", "candidat_gagnant",
                  "bloc_politique", "pct_voix"], Inches(0.3), Inches(4.2)),
    ("INDICATEUR\nSOCIO-ECO", ["id_indic (PK)", "codgeo (FK)", "annee",
                                "population", "emploi", "revenus", "delinquance"], Inches(4.8), Inches(1.7)),
    ("DELTA\nINDICATEUR", ["id_delta (PK)", "codgeo (FK)", "delta_population",
                            "delta_emploi", "delta_revenus", "etat_eco (ENUM)"], Inches(4.8), Inches(4.2)),
    ("PREDICTION", ["id_pred (PK)", "codgeo (FK)", "annee_cible", "scenario",
                    "proba_macron", "proba_lepen", "confidence"], Inches(9.3), Inches(3.1)),
]

for (title, attrs, left, top) in entites:
    w, h = Inches(4.0), Inches(2.1)
    add_bg_rect(slide, BLUE_DARK, left, top, w, Inches(0.4))
    add_bg_rect(slide, BLUE_LIGHT, left, top + Inches(0.4), w, h - Inches(0.4))
    slide.shapes[-1].line.color.rgb = BLUE_MID

    add_text_box(slide, title, left + Inches(0.1), top, w - Inches(0.2), Inches(0.4),
                 font_size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    for j, attr in enumerate(attrs):
        is_pk = 'PK' in attr
        is_fk = 'FK' in attr
        color = ORANGE if is_pk else (BLUE_DARK if is_fk else RGBColor(0x2C, 0x3E, 0x50))
        add_text_box(slide, ('🔑 ' if is_pk else '  ↳ ') + attr,
                     left + Inches(0.1), top + Inches(0.45 + j * 0.27),
                     w - Inches(0.2), Inches(0.27),
                     font_size=9, bold=is_pk, color=color)

add_text_box(slide, "Voir outputs/mcd_schema.png pour le schéma visuel complet",
             Inches(0.3), Inches(6.9), Inches(12), Inches(0.4),
             font_size=10, italic=True, color=GREY, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — MODÈLES TESTÉS ET RÉSULTATS
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK_LAYOUT)
add_bg_rect(slide, RGBColor(0xF8, 0xF9, 0xFA))
add_slide_header(slide, "Comparaison des Modèles ML",
                 "Apprentissage supervisé — Sélection du meilleur algorithme")

modeles = [
    ("Régression Logistique", "72,3%", "72,0%", "72,1%", False),
    ("SVM (RBF)", "75,1%", "74,8%", "75,0%", False),
    ("Random Forest", "80,4%", "80,1%", "80,3%", False),
    ("XGBoost ★", "83,07%", "83,02%", "83,16%", True),
]

headers = ["Modèle", "Accuracy", "F1-Score", "CV (5-fold)"]
col_widths = [Inches(4.0), Inches(2.3), Inches(2.3), Inches(2.3)]
col_starts = [Inches(0.5), Inches(4.7), Inches(7.1), Inches(9.5)]

# En-têtes
for j, (header, left, w) in enumerate(zip(headers, col_starts, col_widths)):
    add_bg_rect(slide, BLUE_DARK, left, Inches(1.6), w, Inches(0.45))
    add_text_box(slide, header, left, Inches(1.62), w, Inches(0.4),
                 font_size=12, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

# Lignes
for i, (nom, acc, f1, cv, is_winner) in enumerate(modeles):
    row_top = Inches(2.1 + i * 0.9)
    bg = GREEN if is_winner else (BLUE_LIGHT if i % 2 == 0 else WHITE)
    for left, w in zip(col_starts, col_widths):
        add_bg_rect(slide, bg, left, row_top, w, Inches(0.8))

    vals = [nom, acc, f1, cv]
    for j, (val, left, w) in enumerate(zip(vals, col_starts, col_widths)):
        txt_color = WHITE if is_winner else RGBColor(0x2C, 0x3E, 0x50)
        add_text_box(slide, val, left, row_top, w, Inches(0.8),
                     font_size=12 if j > 0 else 11, bold=is_winner,
                     color=txt_color, align=PP_ALIGN.CENTER)

add_text_box(slide,
    "★ XGBoost retenu : meilleur compromis accuracy / généralisation (écart CV/test = 0,10%)",
    Inches(0.5), Inches(5.9), Inches(12), Inches(0.5),
    font_size=12, bold=True, color=GREEN, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — RÉSULTATS ET MÉTRIQUES
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK_LAYOUT)
add_bg_rect(slide, RGBColor(0xF8, 0xF9, 0xFA))
add_slide_header(slide, "Résultats du Modèle XGBoost",
                 "Métriques de performance — Jeu de test 20% (25 452 enregistrements)")

metrics = [
    ("83,07%", "Accuracy\n(jeu de test)", BLUE_DARK),
    ("82,99%", "Précision\n(pondérée)", BLUE_MID),
    ("83,07%", "Rappel\n(pondéré)", RGBColor(0x1A, 0x8C, 0x6A)),
    ("83,02%", "F1-Score\n(pondéré)", RGBColor(0x8E, 0x44, 0xAD)),
    ("83,16%\n±0,20%", "Cross-Val.\n5-fold", ORANGE),
]

for i, (val, label, color) in enumerate(metrics):
    add_metric_card(slide, val, label, Inches(0.4 + i * 2.5), Inches(1.7),
                    bg=color, width=Inches(2.2), height=Inches(1.6))

add_text_box(slide, "Prédictions géographiques", Inches(0.5), Inches(3.6),
             Inches(12), Inches(0.4), font_size=14, bold=True, color=BLUE_DARK)

geo_results = [
    ("Niveau Régional", "NOUVELLE-AQUITAINE", "MACRON", "✅", "88,58%", "Stable (76%)"),
    ("Niveau Dép. (11/12)", "Ex. DEUX-SÈVRES", "MACRON", "✅", "92,97%", "Croissance"),
    ("Cas limite", "LOT-ET-GARONNE", "MACRON prédit / LE PEN réel", "❌", "56,80%", "Déclin"),
]

for i, (niveau, zone, pred, ok, conf, etat) in enumerate(geo_results):
    row_top = Inches(4.1 + i * 0.85)
    bg = RGBColor(0xD5, 0xF5, 0xE3) if ok == '✅' else RGBColor(0xFAD7A0)
    add_bg_rect(slide, bg, Inches(0.5), row_top, Inches(12.3), Inches(0.75))
    row_txt = f"{ok}  {niveau} — {zone} : {pred}  (conf. {conf}, état éco. : {etat})"
    add_text_box(slide, row_txt, Inches(0.7), row_top, Inches(12), Inches(0.75),
                 font_size=11, color=RGBColor(0x1B, 0x2B, 0x34))

add_text_box(slide, "Précision départementale globale : 91,7%  (11/12 départements corrects)",
             Inches(0.5), Inches(6.8), Inches(12), Inches(0.4),
             font_size=12, bold=True, color=BLUE_DARK, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — FACTEURS LES PLUS CORRÉLÉS
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK_LAYOUT)
add_bg_rect(slide, RGBColor(0xF8, 0xF9, 0xFA))
add_slide_header(slide, "Indicateur le Plus Corrélé aux Résultats Électoraux",
                 "Feature importance XGBoost — Analyse des facteurs prédictifs")

factors = [
    ("delta_population", 15.95, True, "Variation de la population totale"),
    ("delta_pop_active_1564", 11.96, True, "Dynamique population active 15-64 ans"),
    ("delta_actifs_1564", 10.43, True, "Évolution des personnes actives"),
    ("delta_logements", 9.89, True, "Volume de logements (attractivité)"),
    ("delta_logements_vacants", 9.43, False, "Logements vacants (déclin territorial)"),
    ("delta_revenus_median", 8.21, True, "Revenus médians"),
    ("delta_emploi_salarie", 7.54, True, "Emploi salarié local"),
    ("delta_taux_chomage", 6.89, False, "Taux de chômage (→ vote protestataire)"),
    ("delta_delinquance", 5.32, False, "Taux de délinquance"),
    ("delta_diplome_sup", 4.18, True, "Niveau de diplôme supérieur"),
]

bar_max = 15.95
for i, (feat, imp, positive, desc) in enumerate(factors):
    bar_top = Inches(1.6 + i * 0.56)
    bar_w = Inches(imp / bar_max * 5.0)
    color = GREEN if positive else RED
    add_bg_rect(slide, color, Inches(5.0), bar_top + Inches(0.08), bar_w, Inches(0.38))
    add_text_box(slide, f"{feat} — {desc}",
                 Inches(0.3), bar_top, Inches(4.5), Inches(0.5),
                 font_size=10, color=RGBColor(0x2C, 0x3E, 0x50))
    add_text_box(slide, f"{imp:.2f}%",
                 Inches(5.0) + bar_w + Inches(0.1), bar_top, Inches(1.0), Inches(0.5),
                 font_size=10, bold=True, color=BLUE_DARK)

add_text_box(slide,
    "Conclusion : La variation démographique (delta_population, +15,95%) est le meilleur prédicteur.\n"
    "Zones en croissance → vote Centre/Macron   |   Zones en déclin → vote protestataire/RN",
    Inches(0.3), Inches(7.1), Inches(12.5), Inches(0.6),
    font_size=11, bold=True, color=BLUE_DARK, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — PRÉDICTIONS TEMPORELLES
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK_LAYOUT)
add_bg_rect(slide, RGBColor(0xF8, 0xF9, 0xFA))
add_slide_header(slide, "Scénarios de Prédiction à 1, 2 et 3 Ans",
                 "Projections 2023–2025 — Trois scénarios socio-économiques")

scenarios = [
    ("Scénario Optimiste\n(Croissance économique)", "+1,2% / an",
     ["Gironde maintient >90%", "Lot-et-Garonne remonte à ~65%",
      "Pas de basculement prévu", "Région stable → Macron"], GREEN),
    ("Scénario Central\n(Statu quo)", "-0,3% / an",
     ["Légère érosion du centre", "Lot-et-Garonne reste fragile (~55%)",
      "11/12 dép. restent Macron", "Risque limité"], BLUE_MID),
    ("Scénario Pessimiste\n(Crise économique)", "-2,5% / an",
     ["Recul significatif du centre", "Lot-et-Garonne bascule Le Pen",
      "Creuse et Corrèze fragilisées", "Région passe à ~75%"], RED),
]

for i, (titre, derive, bullets, color) in enumerate(scenarios):
    left = Inches(0.3 + i * 4.35)
    add_bg_rect(slide, color, left, Inches(1.6), Inches(4.1), Inches(0.5))
    add_text_box(slide, titre, left + Inches(0.1), Inches(1.62),
                 Inches(3.9), Inches(0.46), font_size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_bg_rect(slide, BLUE_LIGHT if i == 1 else RGBColor(0xF2, 0xF3, 0xF4),
                left, Inches(2.1), Inches(4.1), Inches(4.0))

    add_text_box(slide, f"Dérive : {derive}", left + Inches(0.2), Inches(2.15),
                 Inches(3.7), Inches(0.4), font_size=12, bold=True, color=color)
    for j, b in enumerate(bullets):
        add_text_box(slide, f"▸  {b}", left + Inches(0.2), Inches(2.6 + j * 0.8),
                     Inches(3.7), Inches(0.75), font_size=11, color=RGBColor(0x2C, 0x3E, 0x50))

add_text_box(slide,
    "Voir outputs/predictions_temporelles.png · predictions_heatmap_scenarios.png · predictions_distribution_y3.png",
    Inches(0.3), Inches(6.9), Inches(12.5), Inches(0.4),
    font_size=9, italic=True, color=GREY, align=PP_ALIGN.CENTER)

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — DÉFINITIONS ANALYTIQUES
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK_LAYOUT)
add_bg_rect(slide, RGBColor(0xF8, 0xF9, 0xFA))
add_slide_header(slide, "Réponses aux Questions Analytiques",
                 "Apprentissage supervisé · Accuracy · Corrélations")

questions = [
    ("Q1 — Donnée la plus corrélée aux élections ?",
     "La variation de population (delta_population, r=+0,52, importance=15,95%) est la plus\n"
     "corrélée. Les zones en croissance démographique votent Centre ; les zones en déclin votent RN.",
     BLUE_DARK),
    ("Q2 — Définir l'apprentissage supervisé",
     "Méthode ML où l'algorithme apprend depuis des exemples étiquetés (X=indicateurs, y=vote).\n"
     "Entraînement sur 80% (101 808 lignes) → généralisation → évaluation sur 20% (25 452 lignes).",
     RGBColor(0x1A, 0x8C, 0x6A)),
    ("Q3 — Définir l'accuracy (degré de précision)",
     "Accuracy = prédictions correctes / total prédictions.\n"
     "Notre modèle : 83,07% sur le jeu de test · 91,7% au niveau départemental · CV 83,16% ±0,20%.\n"
     "Complété par : Précision (82,99%), Rappel (83,07%), F1-Score (83,02%).",
     RGBColor(0x8E, 0x44, 0xAD)),
]

for i, (q, a, color) in enumerate(questions):
    top = Inches(1.6 + i * 1.85)
    add_bg_rect(slide, color, Inches(0.3), top, Inches(12.5), Inches(0.4))
    add_text_box(slide, q, Inches(0.5), top, Inches(12), Inches(0.4),
                 font_size=12, bold=True, color=WHITE)
    add_bg_rect(slide, BLUE_LIGHT if i % 2 == 0 else WHITE,
                Inches(0.3), top + Inches(0.4), Inches(12.5), Inches(1.35))
    add_text_box(slide, a, Inches(0.5), top + Inches(0.42), Inches(12.2), Inches(1.3),
                 font_size=11, color=RGBColor(0x2C, 0x3E, 0x50))

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — LIVRABLES ET ARCHITECTURE TECHNIQUE
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK_LAYOUT)
add_bg_rect(slide, RGBColor(0xF8, 0xF9, 0xFA))
add_slide_header(slide, "Livrables Produits",
                 "Conformité aux exigences Electio-Analytics")

livrables = [
    ("✅", "Dossier de synthèse", "rapport_synthese.md — justification, MCD, démarche, résultats, questions analytiques"),
    ("✅", "Jeu de données nettoyé", "127 260 enregistrements · export CSV · NoSQL MongoDB/TinyDB"),
    ("✅", "Code propre et commenté", "run_prep.py (ETL) · create_bdd.py · ML notebook · Flask app · Chemins portables"),
    ("✅", "Visualisations", "5 types : feature importance, heatmap corrélations, histogrammes, carte, scatter"),
    ("✅", "Modèle prédictif supervisé", "XGBoost · 83% accuracy · train/test split · cross-validation"),
    ("✅", "Prédictions temporelles", "Scénarios 1, 2, 3 ans avec courbes, heat-maps, distributions"),
    ("✅", "MCD (schéma conceptuel)", "mcd.py → mcd_schema.png · 5 entités · relations 1-N"),
    ("✅", "Dashboard interactif", "Flask + HTML/JS · API /api/results · predictions.json"),
    ("✅", "Support de soutenance", "Ce fichier PPTX — 11 slides · synthèse complète"),
]

for i, (icon, titre, desc) in enumerate(livrables):
    top = Inches(1.55 + i * 0.6)
    bg = RGBColor(0xD5, 0xF5, 0xE3) if i % 2 == 0 else WHITE
    add_bg_rect(slide, bg, Inches(0.3), top, Inches(12.5), Inches(0.55))
    add_text_box(slide, f"{icon}  {titre}", Inches(0.5), top, Inches(3.0), Inches(0.55),
                 font_size=11, bold=True, color=BLUE_DARK)
    add_text_box(slide, desc, Inches(3.7), top, Inches(9.0), Inches(0.55),
                 font_size=10, color=RGBColor(0x2C, 0x3E, 0x50))

# ════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — CONCLUSION ET RECOMMANDATIONS
# ════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(BLANK_LAYOUT)
add_bg_rect(slide, BLUE_DARK)
add_bg_rect(slide, ORANGE, 0, Inches(6.8), prs.slide_width, Inches(0.7))

add_text_box(slide, "Conclusion & Recommandations",
             Inches(1), Inches(0.4), Inches(11), Inches(0.6),
             font_size=28, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

conclusions = [
    "Le POC valide l'approche : prédiction électorale possible avec >80% d'accuracy",
    "La variation démographique est le prédicteur le plus puissant (15,95%)",
    "La Nouvelle-Aquitaine reste un bastion Macron à court terme (2023-2025) dans les 3 scénarios",
    "Lot-et-Garonne est la zone à risque de basculement (scénario pessimiste : Le Pen)",
]

for i, c in enumerate(conclusions):
    add_text_box(slide, f"✓  {c}",
                 Inches(1), Inches(1.3 + i * 0.8), Inches(11), Inches(0.7),
                 font_size=15, color=WHITE)

recommandations = [
    "Affiner à l'échelle cantonale pour un meilleur ciblage",
    "Intégrer les données de participation (abstention = signal clé)",
    "Mettre à jour le pipeline dès la publication des données INSEE 2023",
    "Déployer le dashboard Flask en production sur un serveur Linux",
]

add_text_box(slide, "Recommandations :", Inches(1), Inches(4.6), Inches(11), Inches(0.45),
             font_size=16, bold=True, color=ORANGE)
for i, r in enumerate(recommandations):
    add_text_box(slide, f"→  {r}",
                 Inches(1.3), Inches(5.1 + i * 0.45), Inches(11), Inches(0.42),
                 font_size=13, color=BLUE_LIGHT)

add_text_box(slide, "MSPR TPRE813 — Equipe Data Science | Juin 2026 | Electio-Analytics POC",
             Inches(1), Inches(6.85), Inches(11), Inches(0.45),
             font_size=11, color=WHITE, align=PP_ALIGN.CENTER)

# ── Sauvegarde ────────────────────────────────────────────────────────────────
out_path = os.path.join(OUTPUT_DIR, 'soutenance_electio_analytics.pptx')
prs.save(out_path)
print(f"✅ Présentation PowerPoint générée : {out_path}")
print(f"   12 slides : Titre, Contexte, Géographie, ETL, MCD, Modèles,")
print(f"               Résultats, Facteurs, Temporel, Analytique, Livrables, Conclusion")
